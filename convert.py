import os
from pptx import Presentation

def create_presentation():
    # Enter directory of image files here 
    directory = "demo/"
    slidesDir = os.listdir(directory)

    prs = Presentation()

    # Iterates through images in the directory, creates new slide with image
    for slides in slidesDir:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fileName= directory + slides
        try:
            pic = slide.shapes.add_picture(fileName, 0, 0)
        except:
            print(slides + " is not an image file")
            continue

    # Sets size of presentation to the image
    prs.slide_height = pic.height
    prs.slide_width = pic.width
    
    # Enter output directory and name here
    prs.save("output.pptx")


create_presentation()